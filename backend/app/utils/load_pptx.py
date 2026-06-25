from langchain_core.documents import Document

def load_pptx(filepath: str) -> list[Document]:
    from pptx import Presentation
    prs = Presentation(filepath)
    documents = []
    for slide_idx, slide in enumerate(prs.slides):
        slide_text = []
        for shape in slide.shapes:
            text = getattr(shape, "text", None)
            if isinstance(text, str) and text.strip():
                slide_text.append(text.strip())
        content = "\n".join(slide_text)
        if content.strip():
            documents.append(Document(
                page_content=content,
                metadata={
                    "source": filepath,
                    "page": slide_idx + 1,
                    "page_label": str(slide_idx + 1)
                }
            ))
    return documents